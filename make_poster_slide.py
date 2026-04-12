"""
Generate poster_code_section.pptx — full column slide for the
Aviation Tool Accountability poster (System Overview & Code Stack).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY  = RGBColor(0x1a, 0x3a, 0x4a)
TEAL  = RGBColor(0x2e, 0x7d, 0x9e)
GREEN = RGBColor(0x4a, 0x67, 0x41)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x22, 0x22, 0x22)
LGRAY = RGBColor(0xdd, 0xdd, 0xdd)
BGBLUE = RGBColor(0xea, 0xf4, 0xf8)


def add_rect(slide, l, t, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def section_header(slide, text, t, w=5.0):
    shape = add_rect(slide, 0, t, w, 0.28, fill=NAVY)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "  " + text.upper()
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return t + 0.30


def sub_header(slide, text, t, color=TEAL):
    shape = add_rect(slide, 0.18, t, 4.64, 0.24, fill=color)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "  " + text
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return t + 0.26


def bullet_list(slide, items, t, indent=0.36):
    txBox = slide.shapes.add_textbox(Inches(indent), Inches(t), Inches(4.46), Inches(len(items) * 0.22 + 0.05))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(9)
        r.font.color.rgb = DARK
    return t + len(items) * 0.22 + 0.10


def pipeline_row(slide, t):
    boxes = [
        ("capture.py",   "Photo"),
        ("process.py",   "Silhouette"),
        ("vectorize.py", "SVG"),
        ("export.py",    "Laser File"),
    ]
    box_w = 0.88
    gap   = 0.18
    start = 0.18
    for i, (name, sub) in enumerate(boxes):
        x = start + i * (box_w + gap)
        shape = add_rect(slide, x, t, box_w, 0.42, fill=TEAL)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = name + "\n"
        r1.font.size = Pt(8)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = sub
        r2.font.size = Pt(7)
        r2.font.color.rgb = RGBColor(0xCC, 0xE8, 0xF4)
        if i < len(boxes) - 1:
            ax = x + box_w + 0.02
            tb = slide.shapes.add_textbox(Inches(ax), Inches(t + 0.08), Inches(0.14), Inches(0.28))
            p2 = tb.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            r = p2.add_run()
            r.text = "›"
            r.font.size = Pt(14)
            r.font.color.rgb = TEAL
    return t + 0.50


def divider(slide, t):
    add_rect(slide, 0.18, t, 4.64, 0.01, fill=LGRAY)
    return t + 0.10


def pills_row(slide, items, t):
    px, py = 0.18, t
    for pill in items:
        pw = len(pill) * 0.075 + 0.20
        if px + pw > 4.82:
            px = 0.18
            py += 0.30
        shape = add_rect(slide, px, py, pw, 0.23, fill=NAVY)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = pill
        r.font.size = Pt(8)
        r.font.color.rgb = WHITE
        px += pw + 0.08
    return py + 0.33


# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = Inches(5)
prs.slide_height = Inches(11)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

y = 0.0

# ── SYSTEM OVERVIEW ─────────────────────────────────────────────────────────
y = section_header(slide, "System Overview", y)
y = bullet_list(slide, [
    "Backlit light board illuminates tools from below",
    "Camera mounted overhead captures tool silhouette",
    "Raspberry Pi runs the full processing pipeline",
    "ArUco marker provides automatic real-world scale",
    "Outline expanded for foam insert clearance",
    "SVG sent to laser cutter → foam shadow board insert",
], y)

y = divider(slide, y)

# ── PROCESSING PIPELINE ─────────────────────────────────────────────────────
y = section_header(slide, "Processing Pipeline", y)
y = pipeline_row(slide, y)
y += 0.08

y = bullet_list(slide, [
    "capture.py  — Pi camera JPEG capture, locked WB & focus",
    "process.py  — Undistort → detect scale → extract silhouette",
    "vectorize.py — Binary mask → SVG path via vtracer",
    "export.py   — Patch SVG to real-world mm for laser cutter",
    "config.py   — All tunable parameters in one place",
], y)

y = divider(slide, y)

# ── HARDWARE ────────────────────────────────────────────────────────────────
y = section_header(slide, "Hardware", y)
y = bullet_list(slide, [
    "Raspberry Pi",
    "Pi Camera Module",
    "Backlit light board (8\" × 11.5\")",
    "Printed ArUco marker (20mm reference)",
    "Laser cutter",
], y)

y = divider(slide, y)

# ── SOFTWARE ────────────────────────────────────────────────────────────────
y = section_header(slide, "Software", y)
y = bullet_list(slide, [
    "Python 3",
    "Picamera2  — camera capture & control",
    "OpenCV  — distortion correction, ArUco detection, thresholding",
    "NumPy  — image array math",
    "vtracer  — bitmap-to-SVG vectorization",
], y)

y = divider(slide, y)

# ── OUTPUT ──────────────────────────────────────────────────────────────────
y = section_header(slide, "Output", y)
y = bullet_list(slide, [
    "SVG cut file with real-world mm dimensions",
    "Hairline strokes formatted for laser cutter",
    "Laser-cut foam shadow insert for toolbox",
], y)

y = divider(slide, y)

# ── TECH STACK PILLS ────────────────────────────────────────────────────────
y = pills_row(slide, [
    "Python 3", "OpenCV", "Picamera2", "vtracer",
    "NumPy", "Raspberry Pi", "ArUco", "SVG",
], y)

# Trim slide height to content
prs.slide_height = Inches(round(y + 0.15, 2))

prs.save("poster_code_section.pptx")
print(f"Saved: poster_code_section.pptx  (height={round(y+0.15,2)}\")")
